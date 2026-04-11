from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

# Template definitions
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

def add_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    for i, pt in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = pt
        else:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0

# Slide 1
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Global AI Governance Copilot"
subtitle.text = "An Agentic System for Cross-Border AI Compliance\nGenerated System Presentation"

# Slide 2 to 25
slides_data = [
    ("1. The Challenge: Global AI Regulation", ["Global AI laws (e.g., EU AI Act, India Code) are highly fragmented.", "Organizations face severe compliance blind spots when operating across borders.", "Manual review of legal texts is slow, expensive, and prone to error.", "Differing definitions of 'High Risk' create legal conflicts."]),
    ("2. The Solution: AI Governance Copilot", ["An automated, agentic system that analyzes draft AI policy text.", "Compares internal policies against real EU and India legal corpora.", "Instant identification of compliance gaps and conflict signals.", "Scalable foundation for future jurisdictions."]),
    ("3. Core Functional Capabilities", ["Coverage Gap Analysis: Are you missing mandatory clauses?", "Cross-Border Conflict Signals: Do your policies clash with regional laws?", "Policy Option Cards: Choose between Minimal, Moderate, or Strict alignment.", "Automated Reporting: Generating audit-ready PDF reports."]),
    ("4. User Workflow", ["1. Input: User pastes draft AI policy text into the platform.", "2. Process: Agent pipeline parses, classifies, and checks against vector DB.", "3. Dashboard: Interactive UI displays gaps, conflicts, and metrics.", "4. Export: Download final PDF compliance report."]),
    ("5. System Architecture Overview", ["Frontend: Streamlit UI for dynamic, card-based interaction.", "Backend API: FastAPI serving as the robust orchestration layer.", "Data Layer: PostgreSQL (Supabase) storing 5,000+ legal clauses.", "AI Pipeline: Custom NLP agents interacting with FAISS vector search."]),
    ("6. Technology Stack", ["Language: Python 3.x", "API Framework: FastAPI with Uvicorn", "UI Framework: Streamlit with custom dark-mode CSS", "Vector Store: FAISS (cpu) - local indexed search", "Embeddings: sentence-transformers/all-MiniLM-L6-v2"]),
    ("7. Privacy & Security Posture", ["Zero External Model APIs: No draft policies are sent to OpenAI or Anthropic.", "Local Execution: Embeddings generated locally securely via Sentence-Transformers.", "Role-Based Access: Supabase RLS protects compiled legislative data.", "Data Sovereignty: Ensures sensitive unreleased policies remain within the internal network boundary."]),
    ("8. The Legal Data Layer", ["Powered by Supabase (PostgreSQL) for structured storage.", "Contains 5,000+ meticulously extracted clauses from global tech laws.", "Sources: EUR-Lex (EU AI Act) and India Code (IT Act, DPDP).", "Provides grounded, verifiable evidence for every flagged issue."]),
    ("9. Data Pipeline: Building the Corpus", ["Ingestion: Automated extraction of provisions from EUR-Lex PDFs and HTML via PyMuPDF and BeautifulSoup.", "Cleaning: Semantic chunking ensures clauses retain critical context.", "Embedding: Text converted into 384-dimensional dense vectors upon ingestion.", "Database syncing: Automatic population of Supabase and updating of the FAISS index."]),
    ("10. Vector Search & Embeddings", ["Utilizes 'all-MiniLM-L6-v2' for fast, lightweight embedding generation.", "Creates 384-dimensional dense vectors for semantic similarity.", "FAISS IndexFlatIP used for rapid nearest-neighbor retrieval.", "Enables context-aware matching beyond exact keyword searches."]),
    ("11. Agent Pipeline Step 1: Parsing", ["Raw draft policy text is ingested through the UI.", "Input is truncated at 100K characters for memory safety.", "Splits text into discrete, sentence-level actionable clauses.", "Uses spaCy for robust NLP sentence boundary detection, falling back to Regex."]),
    ("12. Agent Pipeline Step 2: Classification", ["Automatically tags each parsed clause using heuristic rules.", "Axes of Classification: Risk Type, Actor Type, Obligation Type.", "Risk Types: High, Unacceptable, Limited, Minimal.", "Actor Types: Provider, Deployer, Importer, Distributor.", "Obligation Types: Transparency, Reporting, Testing, Data Gov."]),
    ("13. Agent Pipeline Step 3: Coverage Checking", ["Evaluates classified clauses against a predefined matrix of requirements.", "Checks essential domains: Biometrics, Incident Reporting, Risk Testing, etc.", "Labels each domain as 'Covered' or a 'Gap'.", "Outputs a detailed grid representing the overall health of the policy."]),
    ("14. Agent Pipeline Step 4: Conflict Detection", ["Cross-references the draft clauses against the EU/India vector database.", "Detects instances where draft text contradicts established law.", "Assigns Severity levels (High, Medium, Low) based on semantic similarity.", "Flags specific regulatory collisions across borders."]),
    ("15. Conflict Detection Deep Dive", ["Similarity Score 0.90+: High Severity Conflict (Red).", "Similarity Score 0.80 - 0.89: Medium Severity Conflict (Amber).", "< 0.80: Low Severity / Contextual mismatch (Green).", "Provides line-by-line evidence comparing draft vs. law excerpt."]),
    ("16. Agent Pipeline Step 5: Recommendations", ["Generates actionable guidance to resolve identified gaps and conflicts.", "Context-aware: Recommends specific language based on the missed areas.", "Highlights priority gaps that must be addressed immediately.", "Powers the three-tiered Policy Alignment Options."]),
    ("17. Policy Alignment Option Cards", ["Minimal Tier: Baseline fixes to achieve absolute minimum compliance.", "Moderate Tier: Industry-standard alignment, balancing safety and speed.", "Strict Tier: Maximum rigor, exceeding baseline laws for high-risk domains.", "Each tier provides custom sample language for immediate insertion."]),
    ("18. Streamlit UI: Premium Experience", ["Custom dark-mode design system with responsive card layouts.", "Interactive Executive Summary with high-level metrics (Clauses, Coverage %).", "Color-coded visual indicators (CSS badges, severity pills).", "Interactive expanders to review how the agent classified raw text."]),
    ("19. System Resiliency & Error Handling", ["Input Validation: Rejects empty inputs and truncates extreme lengths.", "Timeout Hardening: Handles API delays during heavy embedding loads.", "Graceful Degradation: Handles missing database connections safely.", "Clear user feedback via styled error cards in the UI."]),
    ("20. Audit Trails & Export (ReportLab)", ["Automatic generation of timestamped PDF compliance reports.", "Includes coverage matrix, conflict signals, and recommended policies.", "Provides a verifiable artifact for legal teams and board reviews.", "Maintains persistence via database reporting IDs."]),
    ("21. Sample Case: Automated Recruitment", ["Scenario: AI for automated recruitment screening across EU and India.", "Risk Identified: High-Risk biometric verification.", "Action: Triggers mandatory pre-deployment testing and transparency checks.", "Result: Conflict maps ensure both EU AI Act and India DPDP are met."]),
    ("22. Return on Investment (ROI)", ["Time Savings: Reduces multi-jurisdictional review time from weeks to seconds.", "Cost Efficiency: Decreases reliance on costly external legal counsel for preliminary drafts.", "Risk Mitigation: Prevents fines amounting to percentages of global turnover.", "Scalability: Enables rapid market expansion without proportional legal overhead."]),
    ("23. Extensibility & Future Roadmap", ["Keyword classifiers can be seamlessly upgraded to fine-tuned LLMs.", "New regions (e.g., US NIST, APAC rules) easily injected into Supabase.", "Direct integration with CI/CD pipelines to scan 'Policy-as-Code'.", "Multi-lingual embedding models for non-English regulatory texts."]),
    ("24. Business Value & Impact", ["Accelerates compliance audits from weeks to seconds.", "Reduces cross-border legal risk for multinational AI rollouts.", "Provides unified visibility into the organization's regulatory posture.", "Ensures decisions are grounded in actual, current legislation."]),
    ("25. Conclusion", ["AI Governance Copilot transforms manual legal review into an agentic workflow.", "Combines robust software engineering (FastAPI/Streamlit) with modern AI (FAISS).", "Prepares organizations for the complex future of global AI regulation.", "Thank you. Questions?"]),
]

for title_text, points in slides_data:
    add_slide(title_text, points)

prs.save("AI_Governance_Copilot_Presentation.pptx")
print("Presentation successfully updated with 25 slides at AI_Governance_Copilot_Presentation.pptx")
